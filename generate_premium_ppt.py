from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def hex_to_rgb(hex_str):
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

C = {
    'dark': hex_to_rgb("080F1E"),
    'darkCard': hex_to_rgb("111827"),
    'navyMid': hex_to_rgb("0F2340"),
    'cyan': hex_to_rgb("00C8FF"),
    'purple': hex_to_rgb("7C3AED"),
    'gold': hex_to_rgb("F59E0B"),
    'white': hex_to_rgb("FFFFFF"),
    'slate': hex_to_rgb("1E293B"),
    'muted': hex_to_rgb("64748B"),
    'cyanDim': hex_to_rgb("0A4A6B"),
    'red': hex_to_rgb("EF4444"),
    'green': hex_to_rgb("16A34A")
}

def create_premium_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def set_bg(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT):
        shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = align
        return shape

    # --- SLIDE 1: TITLE ---
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s1, C['dark'])
    
    # Right accent
    s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), 0, Inches(4.833), Inches(7.5)).fill.solid()
    s1.shapes[-1].fill.fore_color.rgb = C['navyMid']
    s1.shapes[-1].line.fill.background()
    
    # Vertical cyan line
    s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.48), 0, Inches(0.12), Inches(7.5)).fill.solid()
    s1.shapes[-1].fill.fore_color.rgb = C['cyan']
    s1.shapes[-1].line.fill.background()

    add_text(s1, "MULTI-AGENT AI SYSTEM", 0.55, 0.55, 5.7, 0.32, 12, C['cyan'], bold=True)
    add_text(s1, "AI Debate\n& Hybrid\nIntelligence", 0.5, 1.2, 7.5, 3.5, 64, C['white'], bold=True)
    add_text(s1, "A Framework for Objective Argumentation\nand Structured AI Decision-Making", 0.55, 4.8, 7.5, 0.9, 18, RGBColor(168, 196, 216))

    chips = ["Multi-Agent", "RAG", "XAI", "Hybrid Judge"]
    for i, label in enumerate(chips):
        x_pos = 0.55 + i * 1.8
        rect = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(6.4), Inches(1.6), Inches(0.45))
        rect.fill.solid()
        rect.fill.fore_color.rgb = C['cyanDim']
        rect.line.color.rgb = C['cyan']
        add_text(s1, label, x_pos, 6.4, 1.6, 0.45, 10, C['cyan'], bold=True, align=PP_ALIGN.CENTER)

    # --- SLIDE 2: PROBLEM & SOLUTION ---
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s2, C['dark'])
    add_text(s2, "THE CHALLENGE", 0.5, 0.3, 9, 0.3, 10, C['cyan'], bold=True)
    add_text(s2, "Problem & Solution", 0.5, 0.7, 9, 0.6, 36, C['white'], bold=True)

    # Problem Card
    p_card = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.8))
    p_card.fill.solid()
    p_card.fill.fore_color.rgb = C['darkCard']
    p_card.line.color.rgb = C['red']
    add_text(s2, "THE PROBLEM", 0.8, 2.1, 5, 0.4, 12, C['red'], bold=True)
    problems = ["AI responses are opaque 'black boxes'", "Impossible to trace why one argument wins", "Standard judging is biased & inconsistent", "No structural accountability in reasoning"]
    for i, txt in enumerate(problems):
        add_text(s2, "• " + txt, 0.8, 2.7 + i * 0.8, 5, 0.5, 16, C['white'])

    # Solution Card
    s_card = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.8))
    s_card.fill.solid()
    s_card.fill.fore_color.rgb = C['darkCard']
    s_card.line.color.rgb = C['cyan']
    add_text(s2, "THE SOLUTION", 7.3, 2.1, 5, 0.4, 12, C['cyan'], bold=True)
    solutions = ["Multi-agent debate with structured roles", "Point-Reason-Impact JSON schema enforced", "Transparent 3-layer Hybrid Judge system", "XAI explainability on every decision made"]
    for i, txt in enumerate(solutions):
        add_text(s2, "• " + txt, 7.3, 2.7 + i * 0.8, 5, 0.5, 16, C['white'])

    # --- SLIDE 3: ARCHITECTURE ---
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s3, C['dark'])
    add_text(s3, "ARCHITECTURE", 0.5, 0.3, 9, 0.3, 10, C['cyan'], bold=True)
    add_text(s3, "Multi-Agent Design", 0.5, 0.7, 9, 0.6, 36, C['white'], bold=True)

    # Pro/Con/JSON setup
    def add_agent_card(slide, x, y, title, color, points):
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(3.8), Inches(5.0))
        card.fill.solid()
        card.fill.fore_color.rgb = C['darkCard']
        card.line.color.rgb = color
        add_text(slide, title, x + 0.3, y + 0.3, 3, 0.4, 14, color, bold=True)
        for i, p in enumerate(points):
            add_text(slide, "• " + p, x + 0.3, y + 1.0 + i * 0.8, 3.2, 0.5, 12, C['white'])

    add_agent_card(s3, 0.5, 1.8, "PRO AGENT", C['green'], ["Advocates for the statement", "Logical evidence building", "Structured JSON output", "Persuasive articulation"])
    
    json_card = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.75), Inches(1.8), Inches(3.8), Inches(5.0))
    json_card.fill.solid()
    json_card.fill.fore_color.rgb = hex_to_rgb("0A1628")
    json_card.line.color.rgb = C['cyanDim']
    add_text(s3, "COMMUNICATION PROTOCOL", 4.75, 2.1, 3.8, 0.4, 10, C['cyan'], bold=True, align=PP_ALIGN.CENTER)
    add_text(s3, "JSON Schema Enforcement", 4.75, 2.5, 3.8, 0.4, 14, C['white'], align=PP_ALIGN.CENTER)
    snippet = '{\n  "point": "...",\n  "reason": "...",\n  "impact": "..."\n}'
    add_text(s3, snippet, 5.0, 3.2, 3.3, 2.0, 14, hex_to_rgb("7DD3FC"))

    add_agent_card(s3, 9.0, 1.8, "CON AGENT", C['red'], ["Identifies logical fallacies", "Surfaces risks & weaknesses", "Critical counter-arguments", "Refutes Pro claims"])

    # --- SLIDE 4: RAG ---
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s4, C['dark'])
    add_text(s4, "MEMORY & RETRIEVAL", 0.5, 0.3, 9, 0.3, 10, C['purple'], bold=True)
    add_text(s4, "Retrieval-Augmented Generation (RAG)", 0.5, 0.7, 9, 0.6, 36, C['white'], bold=True)
    
    stats = [
        ("FAISS", "Vector Database", "L2 similarity search for history", C['cyan']),
        ("384-D", "Embeddings", "Sentence-Transformer vectors", C['purple']),
        ("True RAG", "Dynamic Context", "Real-time historical injection", C['gold'])
    ]
    for i, (val, lbl, sub, col) in enumerate(stats):
        x = 0.5 + i * 4.3
        card = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.8), Inches(4.0), Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = C['darkCard']
        card.line.color.rgb = col
        add_text(s4, val, x, 2.5, 4.0, 0.8, 44, C['white'], bold=True, align=PP_ALIGN.CENTER)
        add_text(s4, lbl, x, 3.5, 4.0, 0.4, 18, col, bold=True, align=PP_ALIGN.CENTER)
        add_text(s4, sub, x, 4.2, 4.0, 0.8, 14, C['muted'], align=PP_ALIGN.CENTER)

    # --- SLIDE 5: RULE ENGINE ---
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s5, C['dark'])
    add_text(s5, "HYBRID JUDGE · LAYER 1", 0.5, 0.3, 9, 0.3, 10, C['gold'], bold=True)
    add_text(s5, "Deterministic Rule Engine", 0.5, 0.7, 9, 0.6, 36, C['white'], bold=True)
    
    rules = [
        ("Keyword Analysis", C['cyan'], "Scans for logic connectors"),
        ("Length & Depth", C['gold'], "Reasoning word count check"),
        ("Structural Compliance", C['green'], "JSON field validation")
    ]
    for i, (title, col, sub) in enumerate(rules):
        x = 0.5 + i * 4.3
        card = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.2), Inches(4.0), Inches(4.0))
        card.fill.solid()
        card.fill.fore_color.rgb = C['darkCard']
        card.line.color.rgb = col
        add_text(s5, title, x, 3.0, 4.0, 0.5, 18, col, bold=True, align=PP_ALIGN.CENTER)
        add_text(s5, sub, x, 3.8, 4.0, 0.5, 14, C['white'], align=PP_ALIGN.CENTER)

    # --- SLIDE 6: ML CLASSIFIER ---
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s6, C['dark'])
    add_text(s6, "HYBRID JUDGE · LAYER 2", 0.5, 0.3, 9, 0.3, 10, C['purple'], bold=True)
    add_text(s6, "Logistic Regression Classifier", 0.5, 0.7, 9, 0.6, 36, C['white'], bold=True)
    
    # Large Accuracy
    add_text(s6, "91.7%", 0.5, 2.5, 4.0, 1.5, 84, C['purple'], bold=True, align=PP_ALIGN.CENTER)
    add_text(s6, "Model Accuracy", 0.5, 4.2, 4.0, 0.4, 18, C['white'], align=PP_ALIGN.CENTER)

    # Feature bars
    feats = [("Reasoning Ratio", 88, C['cyan']), ("Keyword Density", 72, C['gold']), ("Structural Score", 91, C['green'])]
    for i, (name, val, col) in enumerate(feats):
        yy = 2.2 + i * 1.2
        add_text(s6, name, 5.0, yy, 3.0, 0.4, 14, C['white'])
        bg_bar = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.0), Inches(yy+0.1), Inches(4.5), Inches(0.3))
        bg_bar.fill.solid()
        bg_bar.fill.fore_color.rgb = C['slate']
        fill_bar = s6.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.0), Inches(yy+0.1), Inches(4.5 * val / 100), Inches(0.3))
        fill_bar.fill.solid()
        fill_bar.fill.fore_color.rgb = col

    # --- SLIDE 7: LLM ---
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s7, C['dark'])
    add_text(s7, "HYBRID JUDGE · LAYER 3", 0.5, 0.3, 9, 0.3, 10, C['cyan'], bold=True)
    add_text(s7, "LLM Semantic Tie-Breaker", 0.5, 0.7, 9, 0.6, 36, C['white'], bold=True)
    
    steps = ["DEADLOCK DETECTED", "PHI-3 LLM CALLED", "VERDICT DELIVERED"]
    for i, txt in enumerate(steps):
        x = 0.5 + i * 4.3
        card = s7.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.5), Inches(4.0), Inches(2.0))
        card.fill.solid()
        card.fill.fore_color.rgb = C['darkCard']
        card.line.color.rgb = C['cyan']
        add_text(s7, txt, x, 3.2, 4.0, 0.6, 18, C['cyan'], bold=True, align=PP_ALIGN.CENTER)

    # --- SLIDE 8: XAI / SHAP ---
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s8, C['dark'])
    add_text(s8, "TRANSPARENCY", 0.5, 0.3, 9, 0.3, 10, C['gold'], bold=True)
    add_text(s8, "Explainable AI (XAI) with SHAP", 0.5, 0.7, 9, 0.6, 36, C['white'], bold=True)
    
    # Waterfall mockup
    add_text(s8, "SHAP Waterfall Plot Breakdown", 0.5, 2.0, 12, 0.5, 18, C['gold'], bold=True)
    shap_data = [("Pro Reason Length", 0.38, C['green']), ("Keyword Density", 0.24, C['green']), ("Con Word Count", -0.15, C['red'])]
    for i, (lbl, val, col) in enumerate(shap_data):
        yy = 3.0 + i * 0.8
        add_text(s8, lbl, 1.0, yy, 4.0, 0.4, 14, C['white'])
        bar_w = abs(val) * 6.0
        bar_x = 6.0 if val > 0 else 6.0 - bar_w
        bar = s8.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(bar_x), Inches(yy), Inches(bar_w), Inches(0.4))
        bar.fill.solid()
        bar.fill.fore_color.rgb = col

    # --- SLIDE 9: EVALUATION ---
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s9, C['dark'])
    add_text(s9, "PERFORMANCE", 0.5, 0.3, 9, 0.3, 10, C['cyan'], bold=True)
    add_text(s9, "Fine-Tuning & Evaluation", 0.5, 0.7, 9, 0.6, 36, C['white'], bold=True)
    
    evals = [("LoRA Fine-Tuning", "1,600+ Samples", C['purple']), ("BLEU Score", "0.74 Accuracy", C['cyan']), ("ROUGE-L", "0.81 Precision", C['gold'])]
    for i, (name, val, col) in enumerate(evals):
        y = 2.0 + i * 1.5
        s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(y), Inches(0.1), Inches(1.2)).fill.solid().fore_color.rgb = col
        add_text(s9, name, 0.8, y + 0.1, 5, 0.4, 18, col, bold=True)
        add_text(s9, val, 0.8, y + 0.5, 5, 0.4, 14, C['white'])

    # --- SLIDE 10: CONCLUSION ---
    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s10, C['dark'])
    
    # Split layout
    s10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.5), 0, Inches(4.833), Inches(7.5)).fill.solid().fore_color.rgb = C['navyMid']
    add_text(s10, "Real-World\nApplications", 0.5, 1.0, 7.5, 2.0, 54, C['cyan'], bold=True)
    
    apps = ["⚖ Legal Research", "🎓 Academic Training", "💼 Executive Decisions"]
    for i, txt in enumerate(apps):
        y = 3.5 + i * 0.8
        chip = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(5.0), Inches(0.6))
        chip.fill.solid().fore_color.rgb = C['darkCard']
        chip.line.color.rgb = C['cyan']
        add_text(s10, txt, 0.7, y + 0.1, 4.5, 0.4, 16, C['white'], bold=True)

    add_text(s10, "Not just intelligent —\nalso objective,\ntransparent &\nauditable.", 9.0, 2.5, 4.0, 3.0, 24, C['white'], bold=True, align=PP_ALIGN.CENTER)

    prs.save('Premium_AI_Debate_Presentation.pptx')
    print("Premium presentation saved successfully as Premium_AI_Debate_Presentation.pptx")

if __name__ == "__main__":
    create_premium_presentation()
